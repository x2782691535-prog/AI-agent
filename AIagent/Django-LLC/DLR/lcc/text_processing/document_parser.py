"""
文档解析器
支持多种文档格式的解析和文本提取
"""

import os
import re
import logging
from typing import Dict, List, Optional, Tuple
from pathlib import Path

# 文档处理库
try:
    import PyPDF2
    from PyPDF2 import PdfReader
except ImportError:
    PyPDF2 = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# OCR功能已移除
OCR_AVAILABLE = False

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    convert_from_path = None
    PDF2IMAGE_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    Image = None
    PIL_AVAILABLE = False

logger = logging.getLogger(__name__)


class DocumentParser:
    """文档解析器"""

    SUPPORTED_FORMATS = {
        '.txt': 'text',
        '.md': 'markdown',
        '.html': 'html',
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.doc': 'doc',
        '.xlsx': 'excel',
        '.xls': 'excel',
        '.pptx': 'powerpoint',
        '.ppt': 'powerpoint'
    }

    def __init__(self):
        """初始化文档解析器"""
        self.encoding_candidates = ['utf-8', 'gbk', 'gb2312', 'utf-16']
    
    def parse_document(self, file_path: str) -> Dict:
        """
        解析文档并提取文本内容
        
        Args:
            file_path: 文档文件路径
            
        Returns:
            Dict: 包含文本内容和元数据的字典
        """
        try:
            file_path = Path(file_path)
            
            if not file_path.exists():
                raise FileNotFoundError(f"文件不存在: {file_path}")
            
            file_ext = file_path.suffix.lower()
            
            if file_ext not in self.SUPPORTED_FORMATS:
                raise ValueError(f"不支持的文件格式: {file_ext}")
            
            # 获取文件基本信息
            file_info = self._get_file_info(file_path)
            
            # 根据文件类型解析内容
            content = self._parse_by_type(file_path, file_ext)
            
            # 合并结果
            result = {
                **file_info,
                'content': content.get('text', ''),
                'metadata': content.get('metadata', {}),
                'success': True,
                'error': None
            }
            
            logger.info(f"成功解析文档: {file_path}")
            return result
            
        except Exception as e:
            logger.error(f"解析文档失败 {file_path}: {e}")
            return {
                'file_path': str(file_path),
                'content': '',
                'metadata': {},
                'success': False,
                'error': str(e)
            }
    
    def _get_file_info(self, file_path: Path) -> Dict:
        """获取文件基本信息"""
        stat = file_path.stat()
        
        return {
            'file_path': str(file_path),
            'file_name': file_path.name,
            'file_size': stat.st_size,
            'file_type': file_path.suffix.lower(),
            'created_time': stat.st_ctime,
            'modified_time': stat.st_mtime
        }
    
    def _parse_by_type(self, file_path: Path, file_ext: str) -> Dict:
        """根据文件类型解析内容"""
        
        if file_ext in ['.txt', '.md']:
            return self._parse_text_file(file_path)
        elif file_ext == '.html':
            return self._parse_html_file(file_path)
        elif file_ext == '.pdf':
            return self._parse_pdf_file(file_path)
        elif file_ext in ['.docx', '.doc']:
            return self._parse_word_file(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            return self._parse_excel_file(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return self._parse_powerpoint_file(file_path)
        else:
            raise ValueError(f"不支持的文件类型: {file_ext}")
    
    def _parse_text_file(self, file_path: Path) -> Dict:
        """解析文本文件"""
        content = self._read_text_with_encoding(file_path)
        
        return {
            'text': content,
            'metadata': {
                'encoding': 'auto-detected',
                'line_count': len(content.splitlines()),
                'char_count': len(content)
            }
        }
    
    def _parse_html_file(self, file_path: Path) -> Dict:
        """解析HTML文件"""
        try:
            from bs4 import BeautifulSoup
            
            content = self._read_text_with_encoding(file_path)
            soup = BeautifulSoup(content, 'html.parser')
            
            # 提取文本内容
            text = soup.get_text(separator='\n', strip=True)
            
            # 提取元数据
            title = soup.find('title')
            title_text = title.get_text() if title else ''
            
            meta_tags = soup.find_all('meta')
            metadata = {
                'title': title_text,
                'meta_tags': {tag.get('name', tag.get('property', '')): tag.get('content', '') 
                             for tag in meta_tags if tag.get('content')}
            }
            
            return {
                'text': text,
                'metadata': metadata
            }
            
        except ImportError:
            # 如果没有BeautifulSoup，使用简单的正则表达式
            content = self._read_text_with_encoding(file_path)
            text = re.sub(r'<[^>]+>', '', content)
            text = re.sub(r'\s+', ' ', text).strip()
            
            return {
                'text': text,
                'metadata': {'parser': 'regex'}
            }
    
    def _parse_pdf_file(self, file_path: Path) -> Dict:
        """解析PDF文件，支持图片PDF检测"""
        if not PyPDF2:
            raise ImportError("需要安装PyPDF2库: pip install PyPDF2")

        try:
            with open(file_path, 'rb') as file:
                reader = PdfReader(file)

                # 首先尝试直接提取文本
                text_parts = []
                text_page_count = 0

                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text().strip()
                    text_parts.append(page_text)
                    if len(page_text) > 50:  # 如果页面有足够的文本
                        text_page_count += 1

                text = '\n'.join(text_parts)
                total_pages = len(reader.pages)

                # 判断是否为图片PDF
                text_ratio = text_page_count / total_pages if total_pages > 0 else 0
                is_image_pdf = text_ratio < 0.3  # 如果少于30%的页面有文本，认为是图片PDF

                logger.info(f"PDF分析: 总页数={total_pages}, 有文本页数={text_page_count}, 文本比例={text_ratio:.2f}")

                if is_image_pdf:
                    logger.warning("检测到图片型PDF（扫描版），无法直接提取文字")
                    # 对于图片PDF，返回说明性文本
                    placeholder_text = f"""这是一个扫描版PDF文档，包含{total_pages}页内容。

由于这是图片格式的PDF，无法直接提取文字内容。
建议的处理方案：

1. 使用在线OCR工具（如白描网页版）进行文字识别
2. 使用专业OCR软件处理
3. 手动输入关键内容

文档信息：
- 文件名: {file_path.name}
- 总页数: {total_pages}
- 有文本的页数: {text_page_count}
- 文本比例: {text_ratio:.1%}

推荐使用在线OCR工具：https://web.baimiaoapp.com/"""

                    text = placeholder_text

                # 提取元数据
                metadata = {
                    'page_count': total_pages,
                    'text_page_count': text_page_count,
                    'text_ratio': text_ratio,
                    'is_image_pdf': is_image_pdf,
                    'requires_manual_processing': is_image_pdf,
                    'pdf_metadata': dict(reader.metadata) if reader.metadata else {}
                }

                return {
                    'text': text,
                    'metadata': metadata
                }

        except Exception as e:
            raise Exception(f"PDF解析失败: {e}")



    def _parse_word_file(self, file_path: Path) -> Dict:
        """解析Word文件"""
        if not DocxDocument:
            raise ImportError("需要安装python-docx库: pip install python-docx")
        
        try:
            doc = DocxDocument(file_path)
            
            # 提取文本
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            text = '\n'.join(text_parts)
            
            # 提取元数据
            metadata = {
                'paragraph_count': len(doc.paragraphs),
                'core_properties': {
                    'title': doc.core_properties.title or '',
                    'author': doc.core_properties.author or '',
                    'subject': doc.core_properties.subject or '',
                    'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                    'modified': str(doc.core_properties.modified) if doc.core_properties.modified else ''
                }
            }
            
            return {
                'text': text,
                'metadata': metadata
            }
            
        except Exception as e:
            raise Exception(f"Word文档解析失败: {e}")
    
    def _parse_excel_file(self, file_path: Path) -> Dict:
        """解析Excel文件"""
        if not openpyxl:
            raise ImportError("需要安装openpyxl库: pip install openpyxl")
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            text_parts = []
            sheet_info = {}
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        sheet_text.append('\t'.join(row_text))
                
                if sheet_text:
                    text_parts.append(f"工作表 {sheet_name}:\n" + '\n'.join(sheet_text))
                    sheet_info[sheet_name] = {
                        'max_row': sheet.max_row,
                        'max_column': sheet.max_column
                    }
            
            text = '\n\n'.join(text_parts)
            
            metadata = {
                'sheet_count': len(workbook.sheetnames),
                'sheet_names': workbook.sheetnames,
                'sheet_info': sheet_info
            }
            
            return {
                'text': text,
                'metadata': metadata
            }
            
        except Exception as e:
            raise Exception(f"Excel文件解析失败: {e}")
    
    def _parse_powerpoint_file(self, file_path: Path) -> Dict:
        """解析PowerPoint文件"""
        if not Presentation:
            raise ImportError("需要安装python-pptx库: pip install python-pptx")
        
        try:
            prs = Presentation(file_path)
            
            text_parts = []
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_parts.append(f"幻灯片 {slide_count}:\n" + '\n'.join(slide_text))
            
            text = '\n\n'.join(text_parts)
            
            metadata = {
                'slide_count': slide_count,
                'core_properties': {
                    'title': prs.core_properties.title or '',
                    'author': prs.core_properties.author or '',
                    'subject': prs.core_properties.subject or '',
                    'created': str(prs.core_properties.created) if prs.core_properties.created else '',
                    'modified': str(prs.core_properties.modified) if prs.core_properties.modified else ''
                }
            }
            
            return {
                'text': text,
                'metadata': metadata
            }
            
        except Exception as e:
            raise Exception(f"PowerPoint文件解析失败: {e}")
    
    def _read_text_with_encoding(self, file_path: Path) -> str:
        """尝试多种编码读取文本文件"""
        for encoding in self.encoding_candidates:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        
        # 如果所有编码都失败，使用错误处理
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            return file.read()
    
    def get_supported_formats(self) -> List[str]:
        """获取支持的文件格式列表"""
        return list(self.SUPPORTED_FORMATS.keys())
    
    def is_supported_format(self, file_path: str) -> bool:
        """检查文件格式是否支持"""
        file_ext = Path(file_path).suffix.lower()
        return file_ext in self.SUPPORTED_FORMATS
